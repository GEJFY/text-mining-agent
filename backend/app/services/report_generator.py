"""レポート生成サービス

テンプレートベースのレポート生成。PPTX/PDF/DOCX/Excel出力。
セクション別データルーティング、実エビデンス参照、セクション間コンテキスト共有。
"""

import json
from datetime import UTC, datetime
from pathlib import Path
from uuid import uuid4

from app.core.logging import get_logger
from app.models.schemas import ReportFormat, ReportRequest, ReportResponse, ReportTemplate
from app.services.llm_orchestrator import LLMOrchestrator, TaskType

logger = get_logger(__name__)

# テンプレート構成定義
TEMPLATE_SECTIONS: dict[ReportTemplate, list[str]] = {
    ReportTemplate.VOC: [
        "エグゼクティブサマリー",
        "感情トレンド分析",
        "クラスター分析結果",
        "主要テーマ別詳細",
        "改善提案",
    ],
    ReportTemplate.AUDIT: [
        "分析概要",
        "主要発見事項",
        "リスク評価",
        "統制上の懸念点",
        "推奨事項",
    ],
    ReportTemplate.COMPLIANCE: [
        "調査概要",
        "時系列分析",
        "キーワード共起分析",
        "リスク分類",
        "結論と提言",
    ],
    ReportTemplate.RISK: [
        "リスク分析概要",
        "リスク分類別集計",
        "ヒートマップ分析",
        "優先対応事項",
        "モニタリング計画",
    ],
}

# セクション→分析タイプのマッピング
SECTION_DATA_MAP: dict[str, list[str]] = {
    # VOC
    "エグゼクティブサマリー": [
        "cluster",
        "cluster_analysis",
        "sentiment",
        "sentiment_analysis",
        "causal_chain_analysis",
        "taxonomy_generation",
    ],
    "感情トレンド分析": ["sentiment", "sentiment_analysis"],
    "クラスター分析結果": ["cluster", "cluster_analysis"],
    "主要テーマ別詳細": [
        "cluster",
        "cluster_analysis",
        "taxonomy_generation",
        "cooccurrence",
        "cooccurrence_analysis",
    ],
    "改善提案": [
        "actionability_scoring",
        "causal_chain_analysis",
        "contradiction_detection",
    ],
    # AUDIT
    "分析概要": [
        "cluster",
        "cluster_analysis",
        "sentiment",
        "sentiment_analysis",
    ],
    "主要発見事項": [
        "cluster_analysis",
        "causal_chain_analysis",
        "contradiction_detection",
        "taxonomy_generation",
    ],
    "リスク評価": [
        "sentiment_analysis",
        "actionability_scoring",
        "causal_chain_analysis",
    ],
    "統制上の懸念点": [
        "contradiction_detection",
        "causal_chain_analysis",
    ],
    "推奨事項": ["actionability_scoring", "causal_chain_analysis"],
    # COMPLIANCE
    "調査概要": [
        "cluster_analysis",
        "sentiment_analysis",
        "taxonomy_generation",
    ],
    "時系列分析": ["sentiment_analysis", "cluster_analysis"],
    "キーワード共起分析": ["cooccurrence", "cooccurrence_analysis"],
    "リスク分類": [
        "taxonomy_generation",
        "actionability_scoring",
        "contradiction_detection",
    ],
    "結論と提言": ["actionability_scoring", "causal_chain_analysis"],
    # RISK
    "リスク分析概要": [
        "cluster_analysis",
        "sentiment_analysis",
        "taxonomy_generation",
    ],
    "リスク分類別集計": ["taxonomy_generation", "cluster_analysis"],
    "ヒートマップ分析": [
        "sentiment_analysis",
        "actionability_scoring",
    ],
    "優先対応事項": ["actionability_scoring", "causal_chain_analysis"],
    "モニタリング計画": [
        "causal_chain_analysis",
        "contradiction_detection",
    ],
}


class ReportGenerator:
    """レポート生成エンジン"""

    def __init__(self, llm: LLMOrchestrator) -> None:
        self.llm = llm
        self.output_dir = Path("reports")
        self.output_dir.mkdir(exist_ok=True)

    async def generate(
        self,
        request: ReportRequest,
        analysis_data: dict,
    ) -> ReportResponse:
        """レポートを生成"""
        report_id = str(uuid4())
        logger.info("report_start", report_id=report_id, template=request.template)

        # セクション構成の取得
        if request.template == ReportTemplate.CUSTOM:
            sections = await self._generate_custom_sections(request.custom_prompt or "")
        else:
            sections = TEMPLATE_SECTIONS.get(request.template, TEMPLATE_SECTIONS[ReportTemplate.VOC])

        # エビデンステキストを事前抽出
        evidence_pool = self._extract_evidence_texts(analysis_data)

        # LLMでセクションコンテンツを生成（セクション間コンテキスト共有）
        report_content = await self._generate_sections(sections, analysis_data, request, evidence_pool)

        # 動的タイトル生成
        title = await self._generate_title(request, analysis_data)

        # 出力形式に応じたファイル生成
        await self._export(report_id, report_content, request.output_format, title=title)

        return ReportResponse(
            report_id=report_id,
            download_url=f"/api/v1/reports/{report_id}/download",
            format=request.output_format,
            generated_at=datetime.now(UTC),
        )

    async def _generate_title(self, request: ReportRequest, analysis_data: dict) -> str:
        """分析内容に基づく動的レポートタイトルを生成"""
        template_name = request.template.value if hasattr(request.template, "value") else str(request.template)
        data_summary = []
        for _atype, adata in analysis_data.items():
            if isinstance(adata, dict):
                result = adata.get("result", adata)
                clusters = result.get("clusters", [])
                if clusters:
                    data_summary.append(f"クラスター{len(clusters)}件")
                dist = result.get("distribution", {})
                if dist:
                    data_summary.append("感情分布あり")
                chains = result.get("chains", [])
                if chains:
                    data_summary.append(f"因果連鎖{len(chains)}件")

        context = ", ".join(data_summary[:5]) if data_summary else "テキストマイニング"
        custom_hint = f" ユーザー指示: {request.custom_prompt[:100]}" if request.custom_prompt else ""

        prompt = f"""以下の分析レポートに最適な日本語タイトル（20文字程度）を1つだけ出力してください。
テンプレート: {template_name}
分析概要: {context}{custom_hint}
タイトルのみ出力（括弧や引用符なし）:"""

        try:
            title = await self.llm.invoke(prompt, TaskType.LABELING, max_tokens=60)
            title = title.strip().strip('"「」').strip()
            if title and len(title) <= 60:
                return title
        except Exception as e:
            logger.warning("title_generation_failed", error=str(e))

        return f"{template_name}分析レポート"

    def _extract_evidence_texts(self, analysis_data: dict) -> list[dict]:
        """分析データからエビデンステキストを収集"""
        evidence = []
        idx = 1

        for atype, adata in analysis_data.items():
            if not isinstance(adata, dict):
                continue
            result = adata.get("result", adata)

            # クラスター分析の代表テキスト
            for cluster in result.get("clusters", []):
                for text in cluster.get("centroid_texts", [])[:2]:
                    evidence.append(
                        {
                            "id": f"E-{idx}",
                            "text": text[:200] if isinstance(text, str) else str(text)[:200],
                            "source": atype,
                            "context": f"クラスター「{cluster.get('title', '')}」",
                        }
                    )
                    idx += 1

            # 感情分析のハイライト
            for h in result.get("highlights", []):
                evidence.append(
                    {
                        "id": f"E-{idx}",
                        "text": h.get("text", "")[:200],
                        "source": atype,
                        "context": f"感情: {h.get('sentiment', '')}",
                    }
                )
                idx += 1

            # 因果連鎖の説明
            for chain in result.get("chains", []):
                chain_str = " → ".join(chain.get("chain", []))
                evidence.append(
                    {
                        "id": f"E-{idx}",
                        "text": chain.get("explanation", chain_str)[:200],
                        "source": atype,
                        "context": f"因果連鎖: {chain_str[:50]}",
                    }
                )
                idx += 1

            # 矛盾検出
            for c in result.get("contradictions", []):
                evidence.append(
                    {
                        "id": f"E-{idx}",
                        "text": f"{c.get('statement_a', '')} vs {c.get('statement_b', '')}",
                        "source": atype,
                        "context": f"矛盾: {c.get('contradiction_type', '')}",
                    }
                )
                idx += 1

            # アクショナビリティ上位
            for item in result.get("items", [])[:5]:
                if item.get("score", 0) >= 0.7:
                    evidence.append(
                        {
                            "id": f"E-{idx}",
                            "text": item.get("text_preview", "")[:200],
                            "source": atype,
                            "context": f"アクション優先度: {item.get('score', 0):.1f}",
                        }
                    )
                    idx += 1

        return evidence[:30]  # 最大30件

    def _format_section_data(self, section_title: str, analysis_data: dict) -> str:
        """セクションに関連する分析データを人間可読な要約に変換"""
        relevant_types = SECTION_DATA_MAP.get(section_title, [])
        parts = []

        for atype in relevant_types:
            adata = analysis_data.get(atype)
            if not adata or not isinstance(adata, dict):
                continue
            result = adata.get("result", adata)

            if atype in ("cluster", "cluster_analysis"):
                clusters = result.get("clusters", [])
                if clusters:
                    lines = [f"[クラスター分析] {len(clusters)}クラスター検出"]
                    for c in clusters[:8]:
                        title = c.get("title", f"Cluster {c.get('cluster_id', '?')}")
                        size = c.get("size", 0)
                        summary = c.get("summary", "")[:100]
                        lines.append(f"  - {title} ({size}件): {summary}")
                    parts.append("\n".join(lines))

            elif atype in ("sentiment", "sentiment_analysis"):
                dist = result.get("distribution", {})
                highlights = result.get("highlights", [])
                if dist or highlights:
                    lines = ["[感情分析]"]
                    if dist:
                        for k, v in dist.items():
                            lines.append(f"  - {k}: {v}")
                    for h in highlights[:5]:
                        lines.append(f"  ★ {h.get('text', '')[:80]} → {h.get('sentiment', '')}")
                    parts.append("\n".join(lines))

            elif atype in ("cooccurrence", "cooccurrence_analysis"):
                nodes = result.get("nodes", [])
                communities = result.get("communities", {})
                if nodes:
                    top5 = sorted(nodes, key=lambda n: n.get("degree_centrality", 0), reverse=True)[:5]
                    lines = [f"[共起ネットワーク] {len(nodes)}ノード"]
                    for n in top5:
                        lines.append(f"  - {n.get('word', '')}: 出現{n.get('frequency', 0)}回")
                    for cid, words in list(communities.items())[:3]:
                        lines.append(f"  コミュニティ{cid}: {', '.join(words[:5])}")
                    parts.append("\n".join(lines))

            elif atype == "causal_chain_analysis":
                chains = result.get("chains", [])
                if chains:
                    lines = [f"[因果連鎖] {len(chains)}チェーン検出"]
                    for c in chains[:5]:
                        arrow = " → ".join(c.get("chain", []))
                        lines.append(f"  - {arrow} (確信度: {c.get('confidence', 0):.1f})")
                    parts.append("\n".join(lines))

            elif atype == "contradiction_detection":
                contradictions = result.get("contradictions", [])
                if contradictions:
                    lines = [f"[矛盾検出] {len(contradictions)}件"]
                    for c in contradictions[:5]:
                        lines.append(
                            f"  - [{c.get('contradiction_type', '')}] "
                            f"{c.get('statement_a', '')[:60]} ⇔ {c.get('statement_b', '')[:60]}"
                        )
                    parts.append("\n".join(lines))

            elif atype == "actionability_scoring":
                items = result.get("items", [])
                if items:
                    lines = [f"[アクショナビリティ] {len(items)}件評価"]
                    top = sorted(items, key=lambda x: x.get("score", 0), reverse=True)[:5]
                    for item in top:
                        lines.append(
                            f"  - [{item.get('category', '')}] スコア{item.get('score', 0):.1f}: "
                            f"{item.get('text_preview', '')[:60]}"
                        )
                    parts.append("\n".join(lines))

            elif atype == "taxonomy_generation":
                root = result.get("root_categories", [])
                if root:
                    lines = [f"[タクソノミー] {len(root)}カテゴリ"]
                    for cat in root[:6]:
                        children = cat.get("children", [])
                        sub = ", ".join(c.get("name", "") for c in children[:3])
                        line = f"  - {cat.get('name', '')}: {cat.get('text_count', 0)}件"
                        if sub:
                            line += f" → {sub}"
                        lines.append(line)
                    parts.append("\n".join(lines))

        if not parts:
            # フォールバック: 全データの簡略表示
            available = [k for k in analysis_data if isinstance(analysis_data[k], dict)]
            if available:
                return f"利用可能な分析: {', '.join(available)}\n（セクション用の詳細データなし）"
            return "分析データなし"

        return "\n\n".join(parts)

    async def _generate_sections(
        self,
        sections: list[str],
        analysis_data: dict,
        request: ReportRequest,
        evidence_pool: list[dict],
    ) -> list[dict]:
        """各セクションのコンテンツをLLMで生成（セクション間コンテキスト共有）"""
        contents = []
        prior_context = ""

        # エビデンスプールのテキスト表現
        evidence_block = ""
        if evidence_pool:
            evidence_lines = []
            for ev in evidence_pool:
                evidence_lines.append(f"[{ev['id']}] ({ev['context']}) {ev['text']}")
            evidence_block = "\n".join(evidence_lines)

        for section_title in sections:
            section_data = self._format_section_data(section_title, analysis_data)

            custom_context = ""
            if request.custom_prompt:
                custom_context = f"\nユーザー指示:\n{request.custom_prompt}\n"

            prompt = f"""テキストマイニング分析レポートの「{section_title}」セクションを作成してください。
{custom_context}
分析データ:
{section_data}

"""
            if evidence_block:
                prompt += f"""エビデンス一覧（[ID]形式で参照可能）:
{evidence_block}

"""
            if prior_context:
                prompt += f"""前セクションまでの要約:
{prior_context}

"""
            prompt += f"""要件:
- ビジネスパーソン向けの明確な文章
- データに基づく具体的な記述
- エビデンスは[E-N]形式で参照を含める
- 前セクションの内容と整合性を保つ
- 200-400字程度

JSON形式:
{{"title": "{section_title}", "content": "...", "evidence_refs": ["E-1", "E-3"]}}"""

            try:
                response = await self.llm.invoke(prompt, TaskType.SUMMARIZATION, max_tokens=1500)
                data = json.loads(response.strip().strip("```json").strip("```"))
                contents.append(data)
                # 次セクション用にコンテキスト蓄積
                content_text = data.get("content", "")
                prior_context += f"【{section_title}】{content_text[:200]}\n"
            except Exception as e:
                logger.warning("section_generation_failed", section=section_title, error=str(e))
                contents.append(
                    {
                        "title": section_title,
                        "content": f"（セクション生成中にエラーが発生しました: {e}）",
                        "evidence_refs": [],
                    }
                )

        return contents

    async def _generate_custom_sections(self, custom_prompt: str) -> list[str]:
        """カスタムプロンプトからセクション構成を生成"""
        prompt = f"""以下の指示に基づいてレポートのセクション構成を設計してください。
指示: {custom_prompt}
JSON配列で5-7セクションのタイトルを出力: ["セクション1", "セクション2", ...]"""

        response = await self.llm.invoke(prompt, TaskType.LABELING, max_tokens=200)
        try:
            return json.loads(response.strip().strip("```json").strip("```"))
        except json.JSONDecodeError:
            return ["概要", "分析結果", "考察", "推奨事項"]

    async def _export(
        self,
        report_id: str,
        contents: list[dict],
        fmt: ReportFormat,
        *,
        title: str = "NexusText AI 分析レポート",
    ) -> Path:
        """各形式でファイルを出力"""
        if fmt == ReportFormat.PPTX:
            return await self._export_pptx(report_id, contents, title=title)
        elif fmt == ReportFormat.PDF:
            return await self._export_pdf(report_id, contents, title=title)
        elif fmt == ReportFormat.DOCX:
            return await self._export_docx(report_id, contents, title=title)
        elif fmt == ReportFormat.EXCEL:
            return await self._export_excel(report_id, contents)
        raise ValueError(f"Unknown format: {fmt}")

    async def _export_pptx(
        self,
        report_id: str,
        contents: list[dict],
        *,
        title: str = "NexusText AI 分析レポート",
    ) -> Path:
        """PowerPoint出力（CJKフォント対応）"""
        from pptx import Presentation
        from pptx.util import Pt

        cjk_font_name = self._find_cjk_font_name() or "Yu Gothic"

        prs = Presentation()

        # タイトルスライド
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = title
        for paragraph in slide.shapes.title.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.name = cjk_font_name
                run.font.size = Pt(28)

        for section in contents:
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = section.get("title", "")
            for paragraph in slide.shapes.title.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.name = cjk_font_name
            body = slide.placeholders[1]
            body.text = section.get("content", "")
            for paragraph in body.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.name = cjk_font_name
                    run.font.size = Pt(14)

        path = self.output_dir / f"{report_id}.pptx"
        prs.save(str(path))
        return path

    async def _export_pdf(
        self,
        report_id: str,
        contents: list[dict],
        *,
        title: str = "NexusText AI 分析レポート",
    ) -> Path:
        """PDF出力（CJKフォント対応）"""
        from reportlab.lib.pagesizes import A4
        from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
        from reportlab.pdfbase import pdfmetrics
        from reportlab.pdfbase.ttfonts import TTFont
        from reportlab.platypus import Paragraph, SimpleDocTemplate, Spacer

        # CJKフォント登録（IPAexGothic → Noto Sans CJK → MSGothic のフォールバック）
        cjk_font = "Helvetica"
        font_candidates = [
            ("/usr/share/fonts/truetype/ipaexfont-gothic/ipaexg.ttf", "IPAexGothic"),
            ("/usr/share/fonts/opentype/ipaexfont-gothic/ipaexg.ttf", "IPAexGothic"),
            ("/usr/share/fonts/opentype/noto/NotoSansCJK-Regular.ttc", "NotoSansCJK"),
            ("/usr/share/fonts/truetype/noto/NotoSansCJK-Regular.ttc", "NotoSansCJK"),
            ("/usr/share/fonts/opentype/noto/NotoSansCJKjp-Regular.otf", "NotoSansCJKjp"),
            ("C:/Windows/Fonts/msgothic.ttc", "MSGothic"),
            ("C:/Windows/Fonts/YuGothM.ttc", "YuGothic"),
        ]
        tried_paths = []
        for font_path, font_name in font_candidates:
            tried_paths.append(font_path)
            if Path(font_path).exists():
                try:
                    pdfmetrics.registerFont(TTFont(font_name, font_path))
                    cjk_font = font_name
                    logger.info("pdf_font_registered", font=font_name, path=font_path)
                    break
                except Exception as e:
                    logger.warning("pdf_font_register_failed", font=font_name, error=str(e))
                    continue

        if cjk_font == "Helvetica":
            logger.warning("pdf_no_cjk_font", tried_paths=tried_paths, fallback="Helvetica")

        path = self.output_dir / f"{report_id}.pdf"
        doc = SimpleDocTemplate(str(path), pagesize=A4)
        styles = getSampleStyleSheet()

        title_style = ParagraphStyle("CJKTitle", parent=styles["Title"], fontName=cjk_font, fontSize=16)
        heading_style = ParagraphStyle("CJKHeading", parent=styles["Heading2"], fontName=cjk_font, fontSize=13)
        normal_style = ParagraphStyle("CJKNormal", parent=styles["Normal"], fontName=cjk_font, fontSize=10, leading=16)

        story = []
        title_escaped = title.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
        story.append(Paragraph(title_escaped, title_style))
        story.append(Spacer(1, 20))

        for section in contents:
            title = section.get("title", "").replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
            content = section.get("content", "").replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
            story.append(Paragraph(title, heading_style))
            story.append(Paragraph(content, normal_style))
            story.append(Spacer(1, 12))

        doc.build(story)
        return path

    async def _export_docx(
        self,
        report_id: str,
        contents: list[dict],
        *,
        title: str = "NexusText AI 分析レポート",
    ) -> Path:
        """Word出力（CJKフォント対応）"""
        from docx import Document
        from docx.shared import Pt

        cjk_font_name = self._find_cjk_font_name()

        doc = Document()
        h = doc.add_heading(title, level=0)
        if cjk_font_name:
            for run in h.runs:
                run.font.name = cjk_font_name
                run.font.size = Pt(18)

        for section in contents:
            h2 = doc.add_heading(section.get("title", ""), level=1)
            if cjk_font_name:
                for run in h2.runs:
                    run.font.name = cjk_font_name
            p = doc.add_paragraph(section.get("content", ""))
            if cjk_font_name:
                for run in p.runs:
                    run.font.name = cjk_font_name

            refs = section.get("evidence_refs", [])
            if refs:
                h3 = doc.add_heading("エビデンス", level=2)
                if cjk_font_name:
                    for run in h3.runs:
                        run.font.name = cjk_font_name
                for ref in refs:
                    bp = doc.add_paragraph(f"• {ref}", style="List Bullet")
                    if cjk_font_name:
                        for run in bp.runs:
                            run.font.name = cjk_font_name

        path = self.output_dir / f"{report_id}.docx"
        doc.save(str(path))
        return path

    def _find_cjk_font_name(self) -> str | None:
        """CJKフォント名を検出"""
        candidates = [
            ("/usr/share/fonts/truetype/ipaexfont-gothic/ipaexg.ttf", "IPAexGothic"),
            ("/usr/share/fonts/opentype/ipaexfont-gothic/ipaexg.ttf", "IPAexGothic"),
            ("/usr/share/fonts/opentype/noto/NotoSansCJK-Regular.ttc", "Noto Sans CJK JP"),
            ("/usr/share/fonts/truetype/noto/NotoSansCJK-Regular.ttc", "Noto Sans CJK JP"),
        ]
        for font_path, font_name in candidates:
            if Path(font_path).exists():
                return font_name
        return None

    async def _export_excel(self, report_id: str, contents: list[dict]) -> Path:
        """Excel出力（CJKフォント対応）"""
        import openpyxl
        from openpyxl.styles import Font

        cjk_font_name = self._find_cjk_font_name() or "Yu Gothic"

        wb = openpyxl.Workbook()
        ws = wb.active
        ws.title = "レポート"

        header_font = Font(name=cjk_font_name, bold=True, size=11)
        cell_font = Font(name=cjk_font_name, size=10)

        ws.append(["セクション", "内容", "エビデンス"])
        for cell in ws[1]:
            cell.font = header_font

        for section in contents:
            ws.append(
                [
                    section.get("title", ""),
                    section.get("content", ""),
                    " | ".join(section.get("evidence_refs", [])),
                ]
            )
            for cell in ws[ws.max_row]:
                cell.font = cell_font

        path = self.output_dir / f"{report_id}.xlsx"
        wb.save(str(path))
        return path
